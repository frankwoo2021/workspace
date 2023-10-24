#########################################################################
# 作者:李宁（蒙娜丽宁），UnityMarvel创始人
#
# 微信公众号：极客起源
#
# B站：https://space.bilibili.com/477001733
#
# Copyright © 2022 Lining. All rights reserved.
#
# 版本: 2.0
#########################################################################

from pptx import Presentation
prs = Presentation("demo.pptx")
for i,slide in enumerate(prs.slides):
    if i == 2:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if text_frame.text != "":
                    print(text_frame.text)



