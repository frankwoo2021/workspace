#########################################################################
# 作者:李宁（蒙娜丽宁），UnityMarvel创始人
#
# 微信公众号：极客起源
#
# B站：https://space.bilibili.com/477001733
#
# Copyright © 2022 Lining. All rights reserved.
#
# 版本: 2.0
#########################################################################

from pptx import Presentation
prs = Presentation("demo.pptx")
from pptx.util import Cm, Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
left = top = width = height = Cm(3)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
tf.text = "这是新添加的文本框"

p = tf.add_paragraph()
p.text = "这是第二段文字，加粗，字号40"
p.font.bold = True
p.font.size = Pt(40)

prs.save("textbox.pptx")