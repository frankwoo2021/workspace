#########################################################################
# 作者:李宁（蒙娜丽宁），UnityMarvel创始人
#
# 微信公众号：极客起源
#
# B站：https://space.bilibili.com/477001733
#
# Copyright © 2022 Lining. All rights reserved.
#
# 版本: 2.0
#########################################################################

from pptx import Presentation
prs = Presentation("demo.pptx")
slide = prs.slides.add_slide(prs.slide_layouts[0])
book = slide.placeholders[0]
author = slide.placeholders[1]

book.text = "《Python从菜鸟到高手》"
author.text = "李宁"
prs.save("demo1.pptx")